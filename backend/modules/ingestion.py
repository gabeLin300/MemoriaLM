import hashlib
import json
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional
from urllib.parse import urlparse

try:
    import requests
except ImportError:
    requests = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

from backend.models.schemas import SourceOut
from backend.services.embeddings import embedding_service
from backend.services.storage import NotebookStore
from backend.services.vector_store import ChromaNotebookStore


ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".txt"}


def _now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def validate_file(path: Path) -> None:
    if path.suffix.lower() not in ALLOWED_EXTENSIONS:
        raise ValueError("Unsupported file type")


def sanitize_filename(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("._")
    return cleaned or "source"


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    if chunk_size <= 0:
        raise ValueError("chunk_size must be > 0")
    if overlap < 0 or overlap >= chunk_size:
        raise ValueError("overlap must be >= 0 and < chunk_size")

    stripped = text.strip()
    if not stripped:
        return []

    chunks: List[str] = []
    i = 0
    while i < len(stripped):
        chunks.append(stripped[i : i + chunk_size])
        i += chunk_size - overlap
    return chunks

# ADDED normalized extracted text and reject empty extractable uploads
def normalize_text(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = "\n".join(line.rstrip() for line in normalized.split("\n"))
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def format_extracted_text(segments: List[Dict[str, str]], source_type: str) -> str:
    cleaned_segments: List[Dict[str, str]] = []
    for seg in segments:
        cleaned = normalize_text(seg.get("text", ""))
        if not cleaned:
            continue
        cleaned_segments.append(
            {
                "location": seg.get("location", "unknown"),
                "text": cleaned,
            }
        )

    if not cleaned_segments:
        return ""

    if source_type == "txt":
        return cleaned_segments[0]["text"]

    blocks: List[str] = []
    for seg in cleaned_segments:
        blocks.append(f"[{seg['location']}]\n{seg['text']}")
    return "\n\n".join(blocks)


def _read_txt(path: Path) -> List[Dict[str, str]]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [{"location": "full text", "text": text}]


def _read_pdf(path: Path) -> List[Dict[str, str]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    segments: List[Dict[str, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            segments.append({"location": f"page {idx}", "text": text})
    return segments


def _read_pptx(path: Path) -> List[Dict[str, str]]:
    from pptx import Presentation

    deck = Presentation(str(path))
    segments: List[Dict[str, str]] = []
    for idx, slide in enumerate(deck.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        combined = "\n".join(texts).strip()
        if combined:
            segments.append({"location": f"slide {idx}", "text": combined})
    return segments


def extract_file_segments(path: Path) -> tuple[str, List[Dict[str, str]]]:
    validate_file(path)
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return "txt", _read_txt(path)
    if suffix == ".pdf":
        return "pdf", _read_pdf(path)
    if suffix == ".pptx":
        return "pptx", _read_pptx(path)
    raise ValueError("Unsupported file type")


def fetch_url_text(url: str, timeout: int = 15) -> tuple[str, str]:
    if requests is None or BeautifulSoup is None:
        raise RuntimeError("URL ingestion dependencies are not installed")
    response = requests.get(url, timeout=timeout, headers={"User-Agent": "MemoriaLM/0.1"})
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")

    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    title = (soup.title.string.strip() if soup.title and soup.title.string else None) or urlparse(url).netloc or "webpage"
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return title, text


def _chunk_segments(
    segments: List[Dict[str, str]],
    *,
    source_id: str,
    source_name: str,
    source_type: str,
    chunk_size: int,
    overlap: int,
) -> List[Dict]:
    chunks: List[Dict] = []
    chunk_counter = 0
    for seg in segments:
        for text_chunk in chunk_text(seg.get("text", ""), chunk_size=chunk_size, overlap=overlap):
            chunk_id = f"{source_id}:{chunk_counter}"
            chunks.append(
                {
                    "chunk_id": chunk_id,
                    "text": text_chunk,
                    "metadata": {
                        "source_id": source_id,
                        "source_name": source_name,
                        "source_type": source_type,
                        "location": seg.get("location", "unknown"),
                        "chunk_index": chunk_counter,
                    },
                }
            )
            chunk_counter += 1
    return chunks


def _source_id_for(content: bytes, source_name: str) -> str:
    digest = hashlib.sha256(content).hexdigest()[:12]
    return f"src_{digest}_{sanitize_filename(source_name)[:32]}"


def _write_source_files(
    store: NotebookStore,
    *,
    user_id: str,
    notebook_id: str,
    source_id: str,
    source_name: str,
    source_type: str,
    raw_bytes: Optional[bytes],
    extracted_text: str,
    chunks: List[Dict],
) -> SourceOut:
    raw_dir = store.files_raw_dir(user_id, notebook_id)
    extracted_dir = store.files_extracted_dir(user_id, notebook_id)

    raw_rel = None
    if raw_bytes is not None:
        raw_ext = Path(source_name).suffix or (".url.txt" if source_type == "url" else "")
        raw_name = f"{source_id}__{sanitize_filename(Path(source_name).stem)}{raw_ext}"
        raw_path = raw_dir / raw_name
        raw_path.write_bytes(raw_bytes)
        raw_rel = str(raw_path.as_posix())

    extracted_name = f"{source_id}.txt"
    extracted_path = extracted_dir / extracted_name
    extracted_path.write_text(extracted_text, encoding="utf-8")

    meta = {
        "source_id": source_id,
        "source_name": source_name,
        "source_type": source_type,
        "enabled": True,
        "raw_path": raw_rel,
        "extracted_path": str(extracted_path.as_posix()),
        "chunk_count": len(chunks),
        "char_count": len(extracted_text),
        "created_at": _now(),
    }
    (extracted_dir / f"{source_id}.meta.json").write_text(json.dumps(meta, indent=2), encoding="utf-8")
    return SourceOut(**meta)


def _index_chunks(store: NotebookStore, user_id: str, notebook_id: str, chunks: List[Dict]) -> None:
    if not chunks:
        return
    embeddings = embedding_service.embed_texts([chunk["text"] for chunk in chunks])
    chroma = ChromaNotebookStore(store.chroma_dir(user_id, notebook_id))
    chroma.upsert_chunks(chunks, embeddings)


def ingest_uploaded_bytes(
    store: NotebookStore,
    *,
    user_id: str,
    notebook_id: str,
    filename: str,
    content: bytes,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> SourceOut:
    filename = sanitize_filename(filename)
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise ValueError("Unsupported file type")

    notebook_dir = store.require_notebook_dir(user_id, notebook_id)
    tmp_path = notebook_dir / f".tmp_{filename}"
    tmp_path.write_bytes(content)
    try:
        source_type, segments = extract_file_segments(tmp_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink()

# ADDED Safeguard for empty documents, upon empty document return error
    extracted_text = format_extracted_text(segments, source_type)
    if not extracted_text:
        raise ValueError("No extractable text found in uploaded file")
    source_id = _source_id_for(content, filename)
    chunks = _chunk_segments(
        segments,
        source_id=source_id,
        source_name=filename,
        source_type=source_type,
        chunk_size=chunk_size,
        overlap=overlap,
    )
    _index_chunks(store, user_id, notebook_id, chunks)
    return _write_source_files(
        store,
        user_id=user_id,
        notebook_id=notebook_id,
        source_id=source_id,
        source_name=filename,
        source_type=source_type,
        raw_bytes=content,
        extracted_text=extracted_text,
        chunks=chunks,
    )


def ingest_url(
    store: NotebookStore,
    *,
    user_id: str,
    notebook_id: str,
    url: str,
    source_name: Optional[str] = None,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> SourceOut:
    title, text = fetch_url_text(url)
    chosen_name = source_name or title or urlparse(url).netloc or "webpage"
    source_name_final = sanitize_filename(chosen_name) + ".url.txt"
    content = text.encode("utf-8")
    source_id = _source_id_for(content, source_name_final)
    segments = [{"location": url, "text": text}]
    chunks = _chunk_segments(
        segments,
        source_id=source_id,
        source_name=source_name_final,
        source_type="url",
        chunk_size=chunk_size,
        overlap=overlap,
    )
    _index_chunks(store, user_id, notebook_id, chunks)
    return _write_source_files(
        store,
        user_id=user_id,
        notebook_id=notebook_id,
        source_id=source_id,
        source_name=source_name_final,
        source_type="url",
        raw_bytes=content,
        extracted_text=text,
        chunks=chunks,
    )


def list_ingested_sources(store: NotebookStore, *, user_id: str, notebook_id: str) -> List[SourceOut]:
    extracted_dir = store.files_extracted_dir(user_id, notebook_id)
    items: List[SourceOut] = []
    for meta_path in sorted(extracted_dir.glob("*.meta.json")):
        try:
            payload = json.loads(meta_path.read_text(encoding="utf-8"))
            if "enabled" not in payload:
                payload["enabled"] = True
            items.append(SourceOut(**payload))
        except Exception:
            continue
    return items


def set_source_enabled(
    store: NotebookStore,
    *,
    user_id: str,
    notebook_id: str,
    source_id: str,
    enabled: bool,
) -> SourceOut:
    extracted_dir = store.files_extracted_dir(user_id, notebook_id)
    meta_path = extracted_dir / f"{source_id}.meta.json"
    if not meta_path.exists():
        raise FileNotFoundError("Source not found")
    payload = json.loads(meta_path.read_text(encoding="utf-8"))
    payload["enabled"] = bool(enabled)
    meta_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    if "enabled" not in payload:
        payload["enabled"] = bool(enabled)
    return SourceOut(**payload)


def enabled_source_ids(store: NotebookStore, *, user_id: str, notebook_id: str) -> set[str]:
    ids: set[str] = set()
    for item in list_ingested_sources(store, user_id=user_id, notebook_id=notebook_id):
        if item.enabled:
            ids.add(item.source_id)
    return ids
